#!/usr/bin/env python3
"""
QED Presentation Generator
Produces a Canva-importable .pptx from the QED tutorial content.
Requires: pip install python-pptx
"""
import sys
sys.stdout.reconfigure(encoding="utf-8", errors="replace")

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ─────────────────────────────────────────────────────────────
# Colour palette  (dark-science theme)
# ─────────────────────────────────────────────────────────────
C_BG         = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy background
C_ACCENT     = RGBColor(0x00, 0xB4, 0xD8)   # electric cyan accent
C_ACCENT2    = RGBColor(0x90, 0xE0, 0xEF)   # light cyan
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_TEXT = RGBColor(0xCA, 0xD3, 0xE0)   # body text
C_GOLD       = RGBColor(0xFF, 0xD1, 0x66)   # highlight / callout
C_DARK_PANEL = RGBColor(0x1A, 0x2E, 0x44)   # card / table bg
C_SUBTLE     = RGBColor(0x44, 0x6B, 0x8C)   # muted accent

# Slide dimensions (16:9 widescreen)
W = Inches(13.33)
H = Inches(7.5)

# ─────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────

def hex_color(r, g, b):
    return RGBColor(r, g, b)


def set_bg(slide, color: RGBColor):
    """Fill slide background with a solid colour."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=C_WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox


def add_paragraph(tf, text, font_size=16, bold=False, color=C_LIGHT_TEXT,
                  align=PP_ALIGN.LEFT, italic=False, space_before=6):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return p


def add_divider(slide, y, color=C_ACCENT, width_pct=0.85):
    left  = W * (1 - width_pct) / 2
    w     = W * width_pct
    shape = slide.shapes.add_shape(1, int(left), int(y), int(w), Pt(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def section_badge(slide, label: str, color=C_ACCENT):
    """Small coloured pill in top-left corner with section number."""
    rect = add_rect(slide, Inches(0.3), Inches(0.22), Inches(1.5), Inches(0.38), color)
    tf = rect.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = C_BG
    run.font.name = "Calibri"


# ─────────────────────────────────────────────────────────────
# Slide builders
# ─────────────────────────────────────────────────────────────

def make_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, C_BG)

    # Decorative top bar
    add_rect(slide, 0, 0, W, Inches(0.12), C_ACCENT)
    # Bottom bar
    add_rect(slide, 0, H - Inches(0.12), W, Inches(0.12), C_ACCENT)

    # Glowing circle decoration (just a circle)
    circle = slide.shapes.add_shape(9, Inches(9.5), Inches(1.0), Inches(5.5), Inches(5.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0x00, 0x3D, 0x5C)
    circle.line.fill.background()

    # Title
    add_textbox(slide, "Quantum Electrodynamics",
                Inches(0.6), Inches(1.4), Inches(8.5), Inches(1.3),
                font_size=44, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    # Subtitle line 1
    add_textbox(slide, "Light and the Speed of Light",
                Inches(0.6), Inches(2.75), Inches(8.5), Inches(0.8),
                font_size=32, bold=False, color=C_ACCENT, align=PP_ALIGN.LEFT)
    # Subtitle line 2
    add_textbox(slide, "A Physicist's Tutorial",
                Inches(0.6), Inches(3.5), Inches(8.5), Inches(0.6),
                font_size=20, bold=False, color=C_LIGHT_TEXT, align=PP_ALIGN.LEFT,
                italic=True)

    add_divider(slide, Inches(4.3), width_pct=0.45)

    add_textbox(slide, "Feynman  ·  Schwinger  ·  Tomonaga  ·  Dirac",
                Inches(0.6), Inches(4.55), Inches(8.0), Inches(0.5),
                font_size=14, color=C_SUBTLE, align=PP_ALIGN.LEFT)
    add_textbox(slide, "Nobel Prize in Physics 1965",
                Inches(0.6), Inches(5.05), Inches(8.0), Inches(0.4),
                font_size=13, color=C_GOLD, align=PP_ALIGN.LEFT, italic=True)


def make_toc_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    add_rect(slide, 0, 0, W, Inches(0.12), C_ACCENT)

    add_textbox(slide, "Presentation Overview",
                Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                font_size=28, bold=True, color=C_WHITE)
    add_divider(slide, Inches(1.15))

    items = [
        ("01", "The Road to QED — Historical Arc"),
        ("02", "Light as a Quantum Field — The Photon"),
        ("03", "Gauge Symmetry U(1) — Why Photons Are Massless"),
        ("04", "The Speed of Light in QED — Why c Is What It Is"),
        ("05", "Feynman's Sum-Over-Paths — How Light Really Travels"),
        ("06", "Virtual vs. Real Photons"),
        ("07", "The Fine Structure Constant α"),
        ("08", "The Quantum Vacuum and Light Propagation"),
        ("09", "QED's Greatest Triumphs — Theory vs. Experiment"),
        ("10", "Modern Frontiers and Open Questions"),
    ]

    col_width = Inches(6.0)
    for i, (num, title) in enumerate(items):
        col = i % 2
        row = i // 2
        left = Inches(0.5) + col * col_width
        top  = Inches(1.4) + row * Inches(1.02)

        add_rect(slide, int(left), int(top), int(col_width - Inches(0.3)),
                 Inches(0.82), C_DARK_PANEL)
        add_textbox(slide, num, int(left) + int(Inches(0.15)), int(top) + int(Inches(0.1)),
                    Inches(0.55), Inches(0.6), font_size=18, bold=True, color=C_ACCENT)
        add_textbox(slide, title, int(left) + int(Inches(0.7)), int(top) + int(Inches(0.1)),
                    int(col_width - Inches(1.1)), Inches(0.6),
                    font_size=14, color=C_LIGHT_TEXT)


def make_section_header(prs, number, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)

    # Full-width accent band
    add_rect(slide, 0, Inches(2.8), W, Inches(2.0), C_DARK_PANEL)
    add_rect(slide, 0, Inches(2.8), Inches(0.12), Inches(2.0), C_ACCENT)

    add_textbox(slide, f"Section {number}",
                Inches(0.5), Inches(1.3), Inches(12), Inches(0.6),
                font_size=20, color=C_SUBTLE, bold=False)
    add_textbox(slide, title,
                Inches(0.5), Inches(2.85), Inches(12), Inches(1.1),
                font_size=36, bold=True, color=C_WHITE)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.5), Inches(3.95), Inches(12), Inches(0.6),
                    font_size=18, color=C_ACCENT2, italic=True)


def make_content_slide(prs, section_num, title, bullets, quote=None, note=None):
    """Standard content slide with bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    add_rect(slide, 0, 0, W, Inches(0.12), C_ACCENT)

    section_badge(slide, f"§{section_num:02d}" if isinstance(section_num, int) else f"§{section_num}")

    add_textbox(slide, title,
                Inches(2.0), Inches(0.18), Inches(10.8), Inches(0.65),
                font_size=24, bold=True, color=C_WHITE)
    add_divider(slide, Inches(0.95))

    # Bullet content area
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.1),
                                     Inches(12.3), Inches(5.2))
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    for b in bullets:
        if b.startswith("##"):
            add_paragraph(tf, b[2:].strip(), font_size=17, bold=True,
                          color=C_ACCENT2, space_before=10 if not first else 0)
        elif b.startswith("  •"):
            add_paragraph(tf, "    " + b[3:].strip(), font_size=15,
                          color=C_LIGHT_TEXT, space_before=3)
        elif b.startswith("•"):
            add_paragraph(tf, "  •  " + b[1:].strip(), font_size=16,
                          color=C_LIGHT_TEXT, space_before=5)
        elif b.startswith("EQ:"):
            eq_box = add_rect(slide, Inches(0.8),
                              Inches(1.1 + tf.paragraphs.__len__() * 0.01),
                              Inches(11.5), Inches(0.5), C_DARK_PANEL)
            add_paragraph(tf, b[3:].strip(), font_size=15, italic=True,
                          color=C_GOLD, space_before=6)
        else:
            add_paragraph(tf, b, font_size=16, color=C_LIGHT_TEXT,
                          space_before=4 if not first else 0)
        first = False

    if quote:
        q_top = Inches(5.5)
        add_rect(slide, Inches(0.4), int(q_top), Inches(12.5), Inches(0.85), C_DARK_PANEL)
        add_rect(slide, Inches(0.4), int(q_top), Inches(0.08), Inches(0.85), C_GOLD)
        add_textbox(slide, f'"{quote}"',
                    Inches(0.65), int(q_top) + int(Pt(6)),
                    Inches(12.1), Inches(0.75),
                    font_size=14, italic=True, color=C_GOLD)

    if note:
        add_textbox(slide, f"▶  {note}",
                    Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.5),
                    font_size=13, color=C_SUBTLE, italic=True)

    return slide


def make_table_slide(prs, section_num, title, headers, rows, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    add_rect(slide, 0, 0, W, Inches(0.12), C_ACCENT)
    section_badge(slide, f"§{section_num:02d}" if isinstance(section_num, int) else f"§{section_num}")

    add_textbox(slide, title,
                Inches(2.0), Inches(0.18), Inches(10.8), Inches(0.65),
                font_size=24, bold=True, color=C_WHITE)
    add_divider(slide, Inches(0.95))

    n_cols = len(headers)
    n_rows = len(rows)
    tbl_left  = Inches(0.5)
    tbl_top   = Inches(1.1)
    tbl_width = Inches(12.3)
    tbl_height = Inches(0.45) * (n_rows + 1)

    table = slide.shapes.add_table(
        n_rows + 1, n_cols, int(tbl_left), int(tbl_top),
        int(tbl_width), int(tbl_height)
    ).table

    col_w = tbl_width // n_cols
    for i in range(n_cols):
        table.columns[i].width = int(col_w)

    def style_cell(cell, text, bg, fg, bold=False, font_size=14):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = fg
        run.font.name = "Calibri"

    for c, h in enumerate(headers):
        style_cell(table.cell(0, c), h, C_ACCENT, C_BG, bold=True, font_size=14)

    for r, row in enumerate(rows):
        bg = C_DARK_PANEL if r % 2 == 0 else C_BG
        for c, val in enumerate(row):
            style_cell(table.cell(r + 1, c), val, bg, C_LIGHT_TEXT, font_size=13)

    if note:
        add_textbox(slide, f"▶  {note}",
                    Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.5),
                    font_size=13, color=C_SUBTLE, italic=True)


def make_equation_slide(prs, section_num, title, equations, explanation=None):
    """Dark slide highlighting key equations."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    add_rect(slide, 0, 0, W, Inches(0.12), C_GOLD)
    section_badge(slide, f"§{section_num:02d}" if isinstance(section_num, int) else f"§{section_num}")

    add_textbox(slide, title,
                Inches(2.0), Inches(0.18), Inches(10.8), Inches(0.65),
                font_size=24, bold=True, color=C_WHITE)
    add_divider(slide, Inches(0.95), color=C_GOLD)

    top = Inches(1.15)
    for eq, desc in equations:
        panel = add_rect(slide, Inches(0.5), int(top), Inches(12.3), Inches(0.72), C_DARK_PANEL)
        add_rect(slide, Inches(0.5), int(top), Inches(0.08), Inches(0.72), C_GOLD)
        add_textbox(slide, eq,
                    Inches(0.75), int(top) + int(Pt(4)),
                    Inches(7.5), Inches(0.62),
                    font_size=17, italic=True, color=C_GOLD, bold=True)
        if desc:
            add_textbox(slide, desc,
                        Inches(8.4), int(top) + int(Pt(8)),
                        Inches(4.3), Inches(0.55),
                        font_size=13, color=C_LIGHT_TEXT)
        top += Inches(0.9)

    if explanation:
        add_textbox(slide, explanation,
                    Inches(0.5), int(top) + int(Inches(0.1)),
                    Inches(12.3), Inches(1.2),
                    font_size=14, color=C_ACCENT2, italic=True)


def make_highlight_slide(prs, section_num, title, stat_pairs, footer=None):
    """Big-number impact slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    add_rect(slide, 0, 0, W, Inches(0.12), C_ACCENT)
    section_badge(slide, f"§{section_num:02d}" if isinstance(section_num, int) else f"§{section_num}")

    add_textbox(slide, title,
                Inches(2.0), Inches(0.18), Inches(10.8), Inches(0.65),
                font_size=24, bold=True, color=C_WHITE)
    add_divider(slide, Inches(0.95))

    n = len(stat_pairs)
    box_w = Inches(12.3) / n
    for i, (stat, label) in enumerate(stat_pairs):
        left = Inches(0.5) + i * box_w
        add_rect(slide, int(left) + int(Inches(0.05)), Inches(1.4),
                 int(box_w) - int(Inches(0.1)), Inches(3.8), C_DARK_PANEL)
        add_textbox(slide, stat,
                    int(left) + int(Inches(0.1)), Inches(1.9),
                    int(box_w) - int(Inches(0.2)), Inches(2.0),
                    font_size=32, bold=True, color=C_ACCENT,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, label,
                    int(left) + int(Inches(0.1)), Inches(3.9),
                    int(box_w) - int(Inches(0.2)), Inches(1.0),
                    font_size=14, color=C_LIGHT_TEXT, align=PP_ALIGN.CENTER)

    if footer:
        add_textbox(slide, footer,
                    Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.9),
                    font_size=14, color=C_GOLD, italic=True, align=PP_ALIGN.CENTER)


def make_closing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    add_rect(slide, 0, 0, W, Inches(0.12), C_ACCENT)
    add_rect(slide, 0, H - Inches(0.12), W, Inches(0.12), C_ACCENT)

    add_textbox(slide, "QED:",
                Inches(0.6), Inches(1.2), Inches(12), Inches(1.0),
                font_size=48, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, "The Most Precisely Tested Theory in Science",
                Inches(0.6), Inches(2.3), Inches(12), Inches(0.8),
                font_size=28, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_divider(slide, Inches(3.3))

    add_textbox(slide,
                '"It has been a mystery ever since it was discovered more than fifty years ago,\n'
                'and all good theoretical physicists put this number up on their wall and worry about it."\n'
                '— Richard P. Feynman, on α ≈ 1/137',
                Inches(1.0), Inches(3.6), Inches(11.3), Inches(1.5),
                font_size=16, italic=True, color=C_GOLD, align=PP_ALIGN.CENTER)

    add_textbox(slide, "Feynman  ·  Schwinger  ·  Tomonaga  ·  Dirac  ·  Maxwell  ·  Einstein",
                Inches(0.6), Inches(5.5), Inches(12), Inches(0.45),
                font_size=13, color=C_SUBTLE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────
# Main build
# ─────────────────────────────────────────────────────────────

def build_presentation(out_path="QED_Presentation.pptx"):
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    # ── 0. Title ──────────────────────────────────────────────
    make_title_slide(prs)

    # ── TOC ───────────────────────────────────────────────────
    make_toc_slide(prs)

    # ══════════════════════════════════════════════════════════
    # SECTION 1 — The Road to QED
    # ══════════════════════════════════════════════════════════
    make_section_header(prs, 1, "The Road to QED", "Three revolutions that had to merge")

    make_table_slide(prs, 1, "Three Revolutions That Had to Merge",
        ["Revolution", "Year", "Key Contribution"],
        [
            ["Maxwell's Electrodynamics", "1865", "c = 1/√(ε₀μ₀) ≈ 3×10⁸ m/s  ·  Light is EM waves"],
            ["Special Relativity (Einstein)", "1905", "Spacetime unification  ·  c invariant for ALL observers"],
            ["Quantum Mechanics", "1900–1926", "Energy quantized  ·  Wave-particle duality"],
        ],
        note="Problem: combining QM + Special Relativity produced negative probabilities and indefinite particle number → a new framework was required."
    )

    make_content_slide(prs, 1, "Dirac's Equation (1928) — The First Seed",
        [
            "•  Dirac sought a relativistic wave equation linear in ∂/∂t and ∇",
            "EQ:  (iγᵘ∂ᵤ − mc/ℏ)ψ = 0",
            "•  Correctly predicted electron spin as a relativistic necessity",
            "•  Predicted antiparticles → positron discovered 1932 ✓",
            "•  Still a single-particle equation — breaks down when pair creation occurs",
            "##  The Key Insight",
            "•  We must promote BOTH the EM field and the electron field to quantum operators",
            "•  Particles are excitations of underlying quantum fields",
        ]
    )

    make_table_slide(prs, 1, "The Founders of QED",
        ["Physicist", "Key Contribution"],
        [
            ["Richard Feynman", "Path integral formulation, Feynman diagrams, pictorial rules"],
            ["Julian Schwinger", "Rigorous operator formalism, renormalization"],
            ["Sin-Itiro Tomonaga", "Independent covariant formulation (Japan, WWII)"],
            ["Freeman Dyson", "Proved all three approaches are mathematically equivalent"],
        ],
        note="Nobel Prize in Physics 1965: Feynman, Schwinger, Tomonaga"
    )

    # ══════════════════════════════════════════════════════════
    # SECTION 2 — The Photon
    # ══════════════════════════════════════════════════════════
    make_section_header(prs, 2, "Light as a Quantum Field", "The photon — precisely defined")

    make_content_slide(prs, 2, "The Electromagnetic Field as a Quantum Object",
        [
            "•  Classical EM: 4-potential Aᵘ = (φ/c, A) satisfying Maxwell's equations",
            "•  QED quantizes this field → photons are quanta",
            "##  Free EM Lagrangian density",
            "EQ:  ℒ_EM = −(1/4) FᵘᵛFᵤᵥ     where  Fᵘᵛ = ∂ᵘAᵛ − ∂ᵛAᵘ",
            "##  Field operator (1D sketch)",
            "EQ:  Â = Σₖ √(ℏ/2ωₖV) [ aₖ e^{i(kx−ωt)} + aₖ† e^{−i(kx−ωt)} ]",
            "  •  aₖ  :  annihilation — removes one photon of momentum ℏk",
            "  •  aₖ† :  creation     — adds one photon of momentum ℏk",
            "  •  Laser beam = coherent state (eigenstate of â)",
        ]
    )

    make_table_slide(prs, 2, "What Is a Photon — Precisely?",
        ["Property", "Value", "Physical Significance"],
        [
            ["Rest mass",   "0",      "Must travel at exactly c · No rest frame exists"],
            ["Spin",        "1",      "Vector boson · Only two transverse polarization states"],
            ["Electric charge", "0",  "Does not self-interact at tree level"],
            ["Helicity",    "±ℏ",     "Left- and right-circular polarization"],
            ["Wavelength",  "E=hf",   "Energy quantized in units hf (Planck 1900)"],
        ],
        note="A laser beam is a coherent state of photons — photon number is maximally uncertain"
    )

    make_content_slide(prs, 2, "The Photon Number Is NOT Always Defined",
        [
            "•  In a coherent state (classical light), photon number N is completely uncertain",
            "EQ:  ΔN · Δφ ≥ 1/2       (number–phase uncertainty relation)",
            "•  Definite phase φ  →  completely indefinite N",
            "•  A Fock state (exactly N photons)  →  completely indefinite phase",
            "##  Consequences",
            "  •  You cannot simultaneously know 'how many photons' and 'what phase'",
            "  •  This is not a measurement limitation — it is a fundamental quantum reality",
            "  •  Single-photon states have been created and detected in lab (quantum optics)",
        ],
        quote="The photon is not a small bright ball — it is a quantum excitation of a field that fills all of space."
    )

    # ══════════════════════════════════════════════════════════
    # SECTION 3 — Gauge Symmetry
    # ══════════════════════════════════════════════════════════
    make_section_header(prs, 3, "Gauge Symmetry U(1)",
                        "The deep reason photons are exactly massless")

    make_equation_slide(prs, 3, "The QED Lagrangian — The Complete Theory",
        [
            ("ℒ_QED = ψ̄(iγᵘDᵤ − mc/ℏ)ψ − (1/4)FᵘᵛFᵤᵥ",
             "Full QED Lagrangian"),
            ("Dᵤ = ∂ᵤ + i(e/ℏc)Aᵤ",
             "Covariant derivative (couples e⁻ to photon)"),
            ("ψ(x) → e^{iα(x)} ψ(x)   and   Aᵤ → Aᵤ − (ℏc/e)∂ᵤα(x)",
             "Local U(1) gauge transformation"),
        ],
        explanation="The Lagrangian is invariant under this transformation for ANY smooth function α(x). "
                    "This local symmetry is the mathematical origin of electromagnetism itself."
    )

    make_content_slide(prs, 3, "Why Gauge Symmetry Forces Zero Photon Mass",
        [
            "•  A photon mass term would look like:",
            "EQ:  ℒ_mass = (1/2)(m_γ c/ℏ)² AᵘAᵤ",
            "•  Under gauge transformation: Aᵤ → Aᵤ + ∂ᵤα",
            "•  So: AᵘAᵤ → (Aᵘ + ∂ᵘα)(Aᵤ + ∂ᵤα) ≠ AᵘAᵤ",
            "•  The mass term BREAKS gauge invariance",
            "##  Conclusion",
            "  •  Gauge symmetry ⟹ m_γ = 0  (exactly, not approximately)",
            "  •  Ward–Takahashi identity protects this at ALL loop orders",
            "  •  Current experimental bound:  m_γ < 10⁻²⁷ eV/c²",
            "  •  This is 22 orders of magnitude smaller than the electron mass",
        ],
        note="Ward–Takahashi: the photon self-energy Πᵘᵛ(q) must always be transverse: qᵤΠᵘᵛ = 0"
    )

    # ══════════════════════════════════════════════════════════
    # SECTION 4 — Speed of Light
    # ══════════════════════════════════════════════════════════
    make_section_header(prs, 4, "The Speed of Light in QED",
                        "Why c is what it is — and what it really means")

    make_content_slide(prs, 4, "c Is the Speed of Causality, Not Just Light",
        [
            "•  c is best understood as the maximum speed of information transfer",
            "•  It is the conversion factor between space and time in Minkowski geometry:",
            "EQ:  ds² = c²dt² − dx² − dy² − dz²",
            "•  Massless particles travel on null geodesics where ds² = 0:",
            "EQ:  ds² = 0  ⟹  c²dt² = |dx|²  ⟹  |v| = c",
            "##  Key insight",
            "  •  This is a geometric consequence of zero rest mass",
            "  •  Not a special property of EM radiation",
            "  •  ANY massless particle (graviton, gluon) travels at exactly c",
            "  •  In Planck units: c = 1 — absorbed into the geometry of spacetime",
        ]
    )

    make_content_slide(prs, 4, "From Maxwell to c — Classical Derivation",
        [
            "•  Maxwell's equations in vacuum give the wave equation:",
            "EQ:  ∇²E = ε₀μ₀ ∂²E/∂t²       ⟹      c = 1/√(ε₀μ₀)",
            "•  In QED: ε₀ and μ₀ arise from the coupling constant e and vacuum structure",
            "•  The fine structure constant encodes EM strength:",
            "EQ:  α = e² / (4πε₀ℏc) ≈ 1/137",
            "##  Running of c in the Quantum Vacuum",
            "  •  Virtual e⁺e⁻ pairs make the vacuum polarizable",
            "  •  QED effective Lagrangian (Euler-Heisenberg) in strong fields:",
            "EQ:  δℒ ~ (α²/90π²)(ℏ³/m⁴c⁵)[4(E²−c²B²)² + 7c²(E·B)²]",
            "  •  Leads to vacuum birefringence in extreme fields (B > B_c ~ 4.4×10⁹ T)",
        ],
        note="Confirmed 2017: vacuum birefringence observed in polarised X-rays from neutron star RX J1856.5−3754"
    )

    make_table_slide(prs, 4, "c in Different Frameworks",
        ["Framework", "Role of c"],
        [
            ["Special Relativity",  "Invariant speed; unifies space and time in Minkowski metric"],
            ["Maxwell Electrodynamics", "Phase velocity of EM waves in vacuum = 1/√(ε₀μ₀)"],
            ["QED",                 "Speed of the massless U(1) gauge boson; set by Lorentz invariance"],
            ["Planck units",        "c = 1 (dimensionless) — geometry absorbs it completely"],
            ["Cosmology",           "Determines light-travel distance and observable universe size"],
        ]
    )

    # ══════════════════════════════════════════════════════════
    # SECTION 5 — Path Integral & Feynman Diagrams
    # ══════════════════════════════════════════════════════════
    make_section_header(prs, 5, "Feynman's Sum-Over-Paths",
                        "How light really travels — all paths simultaneously")

    make_content_slide(prs, 5, "The Path Integral Formulation",
        [
            "•  A quantum particle does not travel a single path",
            "•  It simultaneously explores ALL possible paths",
            "EQ:  ⟨B|A⟩ = ∫ 𝒟[paths] e^{iS[path]/ℏ}",
            "•  Every path contributes a phase factor e^{iS/ℏ}",
            "##  Why Light Appears to Travel in Straight Lines",
            "  •  Near classical path: S varies slowly → phases align → constructive interference",
            "  •  Far from classical path: S varies rapidly → phases cancel → destructive interference",
            "  •  This is the principle of stationary phase: δS = 0 ⟺ classical path",
            "  •  Fermat's principle of least time is recovered as ℏ → 0",
        ],
        quote="Light doesn't really travel in straight lines — it samples all paths. It only appears to because the non-classical paths cancel out. — Feynman"
    )

    make_content_slide(prs, 5, "Feynman Diagram Rules",
        [
            "•  QED solved perturbatively in powers of  α ≈ 1/137",
            "•  Each order represented by a Feynman diagram",
            "##  Diagram elements",
            "  •  Straight lines with arrows  →  electrons / positrons (fermion propagator)",
            "  •  Wavy lines                →  photons (boson propagator)",
            "  •  Each vertex contributes:  −ieγᵘ  and a factor √α to amplitude",
            "##  Why perturbation theory works",
            "EQ:  A = A₀ + A₁(α/π) + A₂(α/π)² + ...      α/π ≈ 0.0023 ≪ 1",
            "  •  Each loop adds a factor of α/π — series converges extremely rapidly",
            "  •  This is why QED is the most precise theory in all of physics",
        ]
    )

    # ══════════════════════════════════════════════════════════
    # SECTION 6 — Virtual vs Real Photons
    # ══════════════════════════════════════════════════════════
    make_section_header(prs, 6, "Virtual vs. Real Photons",
                        "What actually mediates electromagnetic forces")

    make_table_slide(prs, 6, "Real vs. Virtual Photons",
        ["Property", "Real Photon", "Virtual Photon"],
        [
            ["On-shell condition", "pᵘpᵤ = 0  (E = |p|c)", "p² can be anything — spacelike, timelike, null"],
            ["Travel speed", "Exactly c", "Not defined — internal line only"],
            ["Detection", "Yes — carries energy to detector", "No — mathematical propagator"],
            ["Polarization", "Transverse only (2 states)", "All 4 polarizations (gauge-fixed)"],
            ["Physical role", "Real radiation", "Static fields, bound states, forces"],
        ],
        note="The Coulomb 1/r potential emerges from the Fourier transform of the virtual photon propagator 1/q²"
    )

    make_content_slide(prs, 6, "The Coulomb Force — Virtual Photons at Work",
        [
            "•  Electrostatic repulsion = tree-level exchange of one virtual photon",
            "##  Electron scattering diagram",
            "  •  e⁻ ──────●──────────────── e⁻",
            "  •             |  (virtual γ, q² < 0)",
            "  •  e⁻ ──────●──────────────── e⁻",
            "##  Result",
            "EQ:  V(r) = e² / (4πε₀r)       (Coulomb potential)",
            "  •  1/r potential = Fourier transform of propagator 1/q² in 3D",
            "  •  The 'force' is not fundamental — it is emergent from field exchange",
        ],
        quote="Forces are not fundamental — they are the macroscopic shadow of quantum field exchange."
    )

    # ══════════════════════════════════════════════════════════
    # SECTION 7 — Fine Structure Constant
    # ══════════════════════════════════════════════════════════
    make_section_header(prs, 7, "The Fine Structure Constant α",
                        "The number that rules electromagnetism")

    make_equation_slide(prs, 7, "α — Definition and Physical Meaning",
        [
            ("α = e² / (4πε₀ℏc)  ≈  1/137.035 999 084(21)",
             "Dimensionless coupling constant"),
            ("E_n = −α²m_ec²/2n²",
             "Hydrogen energy levels"),
            ("a_e = (g−2)/2 = α/(2π) − 0.328(α/π)² + ...",
             "Electron anomalous magnetic moment"),
            ("σ_Thomson ~ α² / m_e²c⁴",
             "Compton / Thomson cross section"),
        ],
        explanation="α encodes the strength of the coupling between electrons and photons. "
                    "It is the single most important number in atomic physics."
    )

    make_content_slide(prs, 7, "Running of α — It Changes with Energy",
        [
            "•  α is not fixed — it runs with energy scale Q due to vacuum polarization",
            "EQ:  α(Q²) = α(0) / [1 − (α/3π) ln(Q²/m_e²c⁴) + ...]",
            "##  Measured values",
            "  •  Q → 0     (Thomson limit):     α ≈ 1/137.036",
            "  •  Q = m_e c²  (0.511 MeV):        α ≈ 1/137.0",
            "  •  Q = m_Z c²  (91.2 GeV, LEP):    α ≈ 1/128.9",
            "##  Why does α run?",
            "  •  Virtual e⁺e⁻ pairs screen the bare charge at low energies",
            "  •  At higher Q, you probe inside the screening cloud",
            "  •  At very high E: α → O(1) — perturbation theory eventually fails",
        ],
        quote="All good theoretical physicists put this number up on their wall and worry about it. — Feynman"
    )

    # ══════════════════════════════════════════════════════════
    # SECTION 8 — Quantum Vacuum
    # ══════════════════════════════════════════════════════════
    make_section_header(prs, 8, "The Quantum Vacuum",
                        "Empty space is not empty")

    make_content_slide(prs, 8, "What Lives in the Vacuum?",
        [
            "•  The vacuum state |0⟩ is the ground state of the quantum field — NOT nothingness",
            "##  Vacuum contains",
            "  •  Zero-point fluctuations of the EM field:  ⟨0|E²|0⟩ ≠ 0",
            "  •  Virtual particle-antiparticle pairs (constantly appearing/annihilating)",
            "  •  Non-zero vacuum energy density",
            "##  Three measurable consequences",
            "  •  Casimir Effect — force between uncharged conductors",
            "  •  Lamb Shift — atomic energy levels shifted by vacuum fluctuations",
            "  •  Spontaneous Emission — atoms must radiate even in 'empty' space",
        ],
        quote="An atom in empty space is never truly isolated — it is always coupled to vacuum fluctuations of the QED field."
    )

    make_table_slide(prs, 8, "Vacuum Effects — Theory vs. Experiment",
        ["Effect", "QED Prediction", "Measurement", "Status"],
        [
            ["Casimir pressure (d=100nm)", "P = −π²ℏc / 240d⁴", "~1.3×10⁻⁴ N/m²", "Confirmed ~1%"],
            ["Lamb shift (H atom)", "1057.859 MHz", "1057.845 ± 0.10 MHz", "Confirmed"],
            ["Vacuum birefringence", "n∥ ≠ n⊥ in B field", "Observed in neutron star 2017", "Confirmed"],
            ["Schwinger pair production", "Rate ~ exp(−πE_c/E)", "E_c ~ 1.3×10¹⁸ V/m", "Not yet observed"],
        ]
    )

    # ══════════════════════════════════════════════════════════
    # SECTION 9 — Experimental Triumphs
    # ══════════════════════════════════════════════════════════
    make_section_header(prs, 9, "QED's Greatest Triumphs",
                        "The most precisely tested theory in the history of science")

    make_highlight_slide(prs, 9, "Agreement Between QED Theory and Experiment",
        [
            ("12\nsig. figs.", "Electron g−2\nanomalous magnetic moment"),
            ("10\nsig. figs.", "Hydrogen Lamb shift\nenergy splitting"),
            ("1%", "Casimir effect\nvacuum force"),
            ("2017", "Light-by-light\nscattering confirmed\nat LHC (ATLAS)"),
        ],
        footer="Equivalent to measuring the New York–Los Angeles distance to within the width of a human hair"
    )

    make_table_slide(prs, 9, "Electron Anomalous Magnetic Moment — g−2",
        ["Quantity", "Value"],
        [
            ["QED Theory  (aₑ)", "0.001 159 652 181 643 (764)"],
            ["Experiment  (aₑ)", "0.001 159 652 180 73 (28)"],
            ["Agreement",       "12 significant figures ← most precise in all of science"],
            ["Physical meaning", "Electron spin precesses slightly faster than Dirac's prediction due to vacuum"],
        ],
        note="Each additional loop order contributes α/π ≈ 0.0023 — five-loop calculation required for current precision"
    )

    make_table_slide(prs, 9, "All Major QED Tests",
        ["Prediction", "Agreement"],
        [
            ["Electron g−2",              "12 significant figures"],
            ["Muon g−2",                  "10 sig. figs. (4.2σ tension — possible BSM physics!)"],
            ["Hydrogen Lamb shift",       "10 significant figures"],
            ["Compton / Bhabha scattering", "Excellent across all energies"],
            ["Casimir effect",            "~1% precision"],
            ["Light-by-light scattering", "First observed at LHC, 2017"],
            ["Vacuum birefringence",      "Confirmed in neutron star X-ray polarimetry, 2017"],
        ]
    )

    # ══════════════════════════════════════════════════════════
    # SECTION 10 — Open Frontiers
    # ══════════════════════════════════════════════════════════
    make_section_header(prs, 10, "Modern Frontiers",
                        "Where QED meets the unknown")

    make_content_slide(prs, 10, "The Muon g−2 Anomaly — A Crack in the Wall?",
        [
            "•  Muon anomalous magnetic moment shows ~4.2σ deviation from SM prediction",
            "EQ:  Δaμ = a_μ^exp − a_μ^theory ≈ 249 × 10⁻¹¹",
            "•  Fermilab Muon g-2 Experiment (2021–2023) confirms earlier BNL result",
            "##  Possible explanations",
            "  •  New heavy particles coupling to muons (dark photon, Z', leptoquark)",
            "  •  Supersymmetric particles in loop diagrams",
            "  •  Possible QCD hadronic vacuum polarization error (debate ongoing)",
            "##  Status",
            "  •  4.2σ is very suggestive but not yet 5σ 'discovery' threshold",
            "  •  New lattice QCD calculations partially reduce the discrepancy",
            "  •  Resolution expected within 2–3 years",
        ]
    )

    make_content_slide(prs, 10, "Strong-Field QED and Beyond",
        [
            "##  Schwinger Pair Production",
            "•  In fields E ~ E_Schwinger, the vacuum itself becomes unstable",
            "EQ:  E_c = m²c³/(eℏ) ≈ 1.3 × 10¹⁸ V/m",
            "  •  Real e⁺e⁻ pairs spontaneously produced from 'nothing'",
            "  •  Not yet observed — requires ~10× current peak laser intensity",
            "  •  ELI (Extreme Light Infrastructure) aims to approach this regime",
            "##  QED + Gravity",
            "  •  Hawking radiation — black holes emit thermal photons (semi-classical QED)",
            "  •  Unruh effect — accelerating observer sees vacuum as thermal photon bath",
            "  •  Full quantum gravity remains unsolved",
        ]
    )

    make_table_slide(prs, 10, "What QED Does NOT Explain",
        ["Open Question", "Status"],
        [
            ["Why α ≈ 1/137?",             "Unknown — no derivation from first principles"],
            ["Why 3 generations of fermions?", "Unknown"],
            ["Why is e (electron charge) its specific value?", "Unknown"],
            ["Magnetic monopoles?",         "Not observed; predicted by grand unified theories"],
            ["Is the photon truly massless to all orders?", "Yes — protected by gauge symmetry"],
            ["Muon g−2 anomaly?",           "Under active investigation — possible BSM signal"],
        ]
    )

    # ── Key Equations Reference ───────────────────────────────
    make_equation_slide(prs, "∑", "Key Equations — Quick Reference",
        [
            ("ℒ_QED = ψ̄(iγᵘDᵤ − mc/ℏ)ψ − (1/4)FᵘᵛFᵤᵥ",  "Full QED Lagrangian"),
            ("α = e²/(4πε₀ℏc) ≈ 1/137.036",                "Fine structure constant"),
            ("E = |p|c   (m_γ = 0)",                         "Photon dispersion — massless"),
            ("g/2 = 1 + α/(2π) − 0.3285(α/π)² + ...",       "Electron g-factor (QED)"),
            ("P_Casimir = −π²ℏc / (240d⁴)",                  "Casimir pressure"),
            ("E_c = m²c³/(eℏ) ≈ 1.3 × 10¹⁸ V/m",           "Schwinger critical field"),
        ]
    )

    # ── Closing ───────────────────────────────────────────────
    make_closing_slide(prs)

    prs.save(out_path)
    print(f"Saved: {out_path}  ({prs.slides.__len__()} slides)")


if __name__ == "__main__":
    out = sys.argv[1] if len(sys.argv) > 1 else "QED_Presentation.pptx"
    build_presentation(out)
